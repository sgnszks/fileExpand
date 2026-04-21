"""
集成测试

覆盖完整的文件膨胀流程:
- XLSX 膨胀与验证
- DOCX 膨胀与验证
- PPTX 膨胀与验证
- PDF 膨胀与验证

每个测试创建一个最小有效的目标格式文件，
执行膨胀处理，然后验证输出文件的有效性和体积。
"""
import os
import io
import tempfile

import pytest


class TestXlsxIntegration:
    """XLSX 完整流程集成测试。"""

    def test_xlsx_expand_2x(self, temp_dir):
        """测试 XLSX 文件 2 倍膨胀。"""
        import openpyxl
        from app.processors.xlsx_processor import XlsxProcessor
        from app.validators.file_validator import FileValidator

        input_path = os.path.join(temp_dir, 'test.xlsx')
        output_path = os.path.join(temp_dir, 'test_expanded.xlsx')

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = 'TestSheet'
        for i in range(1, 51):
            ws.cell(row=i, column=1, value=f'数据行 {i}')
            ws.cell(row=i, column=2, value=i * 100)
        wb.save(input_path)

        processor = XlsxProcessor()
        result = processor.expand(input_path, output_path, 2.0)

        assert result.success is True
        assert result.output_size > result.original_size

        validator = FileValidator()
        validation = validator.validate(output_path, '.xlsx')
        assert validation['valid'] is True

        wb2 = openpyxl.load_workbook(output_path, read_only=True)
        assert 'TestSheet' in wb2.sheetnames
        wb2.close()

    def test_xlsx_expand_5x(self, temp_dir):
        """测试 XLSX 文件 5 倍膨胀（使用组合策略）。"""
        import openpyxl
        from app.processors.xlsx_processor import XlsxProcessor
        from app.validators.file_validator import FileValidator

        input_path = os.path.join(temp_dir, 'test.xlsx')
        output_path = os.path.join(temp_dir, 'test_expanded_5x.xlsx')

        wb = openpyxl.Workbook()
        ws = wb.active
        ws['A1'] = 'Hello'
        wb.save(input_path)

        processor = XlsxProcessor()
        result = processor.expand(input_path, output_path, 5.0)

        assert result.success is True

        validator = FileValidator()
        validation = validator.validate(output_path, '.xlsx')
        assert validation['valid'] is True


class TestDocxIntegration:
    """DOCX 完整流程集成测试。"""

    def test_docx_expand_2x(self, temp_dir):
        """测试 DOCX 文件 2 倍膨胀。"""
        from docx import Document
        from app.processors.docx_processor import DocxProcessor
        from app.validators.file_validator import FileValidator

        input_path = os.path.join(temp_dir, 'test.docx')
        output_path = os.path.join(temp_dir, 'test_expanded.docx')

        doc = Document()
        doc.add_paragraph('这是测试段落第一行')
        doc.add_paragraph('This is the second paragraph.')
        doc.save(input_path)

        processor = DocxProcessor()
        result = processor.expand(input_path, output_path, 2.0)

        assert result.success is True

        validator = FileValidator()
        validation = validator.validate(output_path, '.docx')
        assert validation['valid'] is True

        doc2 = Document(output_path)
        assert len(doc2.paragraphs) >= 2


class TestPptxIntegration:
    """PPTX 完整流程集成测试。"""

    def test_pptx_expand_2x(self, temp_dir):
        """测试 PPTX 文件 2 倍膨胀。"""
        from pptx import Presentation
        from app.processors.pptx_processor import PptxProcessor
        from app.validators.file_validator import FileValidator

        input_path = os.path.join(temp_dir, 'test.pptx')
        output_path = os.path.join(temp_dir, 'test_expanded.pptx')

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = '测试幻灯片'
        prs.save(input_path)

        processor = PptxProcessor()
        result = processor.expand(input_path, output_path, 2.0)

        assert result.success is True

        validator = FileValidator()
        validation = validator.validate(output_path, '.pptx')
        assert validation['valid'] is True

        prs2 = Presentation(output_path)
        assert len(prs2.slides) == 1


class TestPdfIntegration:
    """PDF 完整流程集成测试。"""

    def test_pdf_expand_2x(self, temp_dir):
        """测试 PDF 文件 2 倍膨胀。"""
        import pikepdf
        from app.processors.pdf_processor import PdfProcessor
        from app.validators.file_validator import FileValidator

        input_path = os.path.join(temp_dir, 'test.pdf')
        output_path = os.path.join(temp_dir, 'test_expanded.pdf')

        pdf = pikepdf.Pdf.new()
        page = pikepdf.Dictionary(
            Type=pikepdf.Name('/Page'),
            MediaBox=[0, 0, 612, 792],
            Contents=pdf.make_indirect(
                pikepdf.Stream(pdf, b'BT /F1 12 Tf 100 700 Td (Test Page) Tj ET')
            ),
            Resources=pikepdf.Dictionary(
                Font=pikepdf.Dictionary(
                    F1=pikepdf.Dictionary(
                        Type=pikepdf.Name('/Font'),
                        Subtype=pikepdf.Name('/Type1'),
                        BaseFont=pikepdf.Name('/Helvetica'),
                    )
                )
            )
        )
        pdf.pages.append(page)
        pdf.save(input_path)

        processor = PdfProcessor()
        result = processor.expand(input_path, output_path, 2.0)

        assert result.success is True

        validator = FileValidator()
        validation = validator.validate(output_path, '.pdf')
        assert validation['valid'] is True

        pdf2 = pikepdf.open(output_path)
        assert len(pdf2.pages) == 1
        pdf2.close()


class TestBinaryFormatRejection:
    """旧版二进制格式拒绝测试。"""

    def test_doc_rejected(self, temp_dir):
        """确认 DOC 格式被正确拒绝。"""
        from app.processors.binary_processor import BinaryOfficeProcessor

        input_path = os.path.join(temp_dir, 'test.doc')
        output_path = os.path.join(temp_dir, 'test_expanded.doc')
        with open(input_path, 'wb') as f:
            f.write(b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1' + b'\x00' * 512)

        processor = BinaryOfficeProcessor()
        result = processor.expand(input_path, output_path, 2.0)

        assert result.success is False
        assert '不支持' in result.error_message

    def test_ppt_rejected(self, temp_dir):
        """确认 PPT 格式被正确拒绝。"""
        from app.processors.binary_processor import BinaryOfficeProcessor

        input_path = os.path.join(temp_dir, 'test.ppt')
        output_path = os.path.join(temp_dir, 'test_expanded.ppt')
        with open(input_path, 'wb') as f:
            f.write(b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1' + b'\x00' * 512)

        processor = BinaryOfficeProcessor()
        result = processor.expand(input_path, output_path, 2.0)

        assert result.success is False
        assert '不支持' in result.error_message
